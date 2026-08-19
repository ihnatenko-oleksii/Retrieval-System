import json
import os

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from app.core.models import Document, DocumentMetadata


class BaseLoader:
    def load(self, file_path: str) -> list[Document]:
        raise NotImplementedError

    def _create_document(self, file_path: str, content: str) -> Document:
        file_name = os.path.basename(file_path)
        extension = os.path.splitext(file_name)[1].lower()
        metadata = DocumentMetadata(
            file_path=file_path, file_name=file_name, extension=extension, loader_type=self.__class__.__name__
        )
        return Document(content=content.strip(), metadata=metadata)


class TextLoader(BaseLoader):
    def load(self, file_path: str) -> list[Document]:
        try:
            with open(file_path, encoding="utf-8") as f:
                content = f.read()
            return [self._create_document(file_path, content)]
        except UnicodeDecodeError:
            with open(file_path, encoding="latin-1") as f:
                content = f.read()
            return [self._create_document(file_path, content)]


class PDFLoader(BaseLoader):
    def load(self, file_path: str) -> list[Document]:
        documents = []
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            content = page.extract_text()
            if content:
                doc = self._create_document(file_path, content)
                doc.metadata.page_number = i + 1
                documents.append(doc)
        return documents


class DocxLoader(BaseLoader):
    def load(self, file_path: str) -> list[Document]:
        doc = DocxDocument(file_path)
        content = "\n".join([para.text for para in doc.paragraphs if para.text])
        if content:
            return [self._create_document(file_path, content)]
        return []


class PptxLoader(BaseLoader):
    def load(self, file_path: str) -> list[Document]:
        prs = Presentation(file_path)
        documents = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text_runs.append(paragraph.text)
            content = "\n".join(text_runs)
            if content.strip():
                doc = self._create_document(file_path, content)
                doc.metadata.page_number = i + 1
                documents.append(doc)
        return documents


class CSVLoader(BaseLoader):
    def load(self, file_path: str) -> list[Document]:
        df = pd.read_csv(file_path)
        content = df.to_string(index=False)
        return [self._create_document(file_path, content)]


class JSONLoader(BaseLoader):
    def load(self, file_path: str) -> list[Document]:
        with open(file_path, encoding="utf-8") as f:
            data = json.load(f)
        content = json.dumps(data, indent=2)
        return [self._create_document(file_path, content)]


def get_loader(file_path: str) -> BaseLoader | None:
    ext = os.path.splitext(file_path)[1].lower()
    text_extensions = {".txt", ".md", ".py", ".java", ".js", ".ts", ".yml", ".yaml", ".xml", ".properties"}
    if ext in text_extensions:
        return TextLoader()
    elif ext == ".pdf":
        return PDFLoader()
    elif ext == ".docx":
        return DocxLoader()
    elif ext == ".pptx":
        return PptxLoader()
    elif ext == ".csv":
        return CSVLoader()
    elif ext == ".json":
        return JSONLoader()
    return None
